"""Build the Ayni pitch deck via python-pptx.

Lightning-themed dark deck, 16:9 widescreen, 8 slides, big type, minimal text.
Output: ayni-pitch/Ayni-pitch-deck.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Palette (Lightning Yellow + Midnight) ---
BG       = RGBColor(0x0B, 0x10, 0x20)  # deep midnight
CARD     = RGBColor(0x18, 0x21, 0x38)  # slate card
CARD_HI  = RGBColor(0x22, 0x2D, 0x4A)  # lighter slate
ACCENT   = RGBColor(0xF7, 0x93, 0x1A)  # bitcoin / lightning orange
ACCENT_2 = RGBColor(0xFF, 0xD5, 0x80)  # warm sand
TEXT     = RGBColor(0xEC, 0xEC, 0xF0)  # off-white
MUTED    = RGBColor(0x8B, 0x93, 0xA7)  # slate-gray
DANGER   = RGBColor(0xE5, 0x5A, 0x4F)  # muted red for "negative-sum"

HEADER_FONT = "Calibri"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, *, line=None, line_w=Pt(1), corner=False, corner_radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    if corner and corner_radius is not None:
        # adj is fraction of half-min-dim, 0..0.5
        shp.adjustments[0] = corner_radius
    # remove default shadow
    sp = shp.fill._xPr  # type: ignore[attr-defined]
    return shp


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, italic=False, color=None,
             align="left", font=BODY_FONT, anchor="top",
             line_spacing=1.1):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tb


def add_runs(slide, runs_spec, left, top, width, height, *,
             align="left", anchor="top", line_spacing=1.15):
    """runs_spec: list of dicts {text, size, bold, italic, color, font, newline_after}"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p = tf.paragraphs[0]
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    p.line_spacing = line_spacing
    for spec in runs_spec:
        run = p.add_run()
        run.text = spec["text"]
        run.font.name = spec.get("font", BODY_FONT)
        run.font.size = Pt(spec.get("size", 18))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        if "color" in spec:
            run.font.color.rgb = spec["color"]
        if spec.get("newline_after"):
            p = tf.add_paragraph()
            p.alignment = align_map.get(align, PP_ALIGN.LEFT)
            p.line_spacing = line_spacing
    return tb


def add_accent_bar(slide, left, top, height, *, width=Inches(0.07), color=ACCENT):
    return add_rect(slide, left, top, width, height, color)


def slide_blank(prs):
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)


# ---------------- BUILD ----------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# === Slide 1: Title ===
s = slide_blank(prs)
set_bg(s, BG)

# subtle accent corner block top-right
add_rect(s, Inches(11.8), Inches(0), Inches(1.6), Inches(0.5), ACCENT)

# Massive "ayni"
add_text(s, "ayni", Inches(0), Inches(1.4), SLIDE_W, Inches(3.2),
         size=180, bold=True, color=ACCENT, align="center", font=HEADER_FONT,
         anchor="middle", line_spacing=1.0)

# Subtitle
add_text(s, "Programmable reciprocity for the agent economy.",
         Inches(0), Inches(4.55), SLIDE_W, Inches(0.6),
         size=24, bold=False, color=TEXT, align="center", font=HEADER_FONT)

# Quechua etymology — single line, only on this slide
add_text(s, "Quechua  ·  today I help you, tomorrow you help me",
         Inches(0), Inches(5.15), SLIDE_W, Inches(0.4),
         size=14, italic=True, color=MUTED, align="center", font=BODY_FONT)

# Team
add_text(s, "Dennis Vivas   ·   Cindy Rojas   ·   Miluska Pajuelo   ·   Jhoselyn Pajuelo",
         Inches(0), Inches(6.45), SLIDE_W, Inches(0.4),
         size=14, color=TEXT, align="center", font=BODY_FONT)

# Footer
add_text(s, "Lima, Peru   ·   Hack Nation 5   ·   Spiral Challenge 02   ·   ayniw.com",
         Inches(0), Inches(6.9), SLIDE_W, Inches(0.4),
         size=11, color=MUTED, align="center", font=BODY_FONT)


# === Slide 2: What we do ===
s = slide_blank(prs)
set_bg(s, BG)

add_accent_bar(s, Inches(0.6), Inches(0.7), Inches(0.5))
add_text(s, "WHAT WE DO", Inches(0.85), Inches(0.7), Inches(4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
         line_spacing=1.0)

# Two big statements
add_text(s,
         "AI agents pay tribes of human experts via Lightning Network.",
         Inches(0.6), Inches(1.9), Inches(12.13), Inches(1.6),
         size=44, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.15)

add_runs(s, [
    {"text": "Every sub-cent payment ", "size": 36, "bold": True, "color": TEXT, "font": HEADER_FONT},
    {"text": "auto-splits", "size": 36, "bold": True, "color": ACCENT, "font": HEADER_FONT},
    {"text": " across all the contributors — in the same second.", "size": 36, "bold": True, "color": TEXT, "font": HEADER_FONT},
], Inches(0.6), Inches(4.0), Inches(12.13), Inches(2.4), line_spacing=1.2)

# Footer
add_text(s, "ayniw.com", Inches(0.6), Inches(6.9), Inches(12.13), Inches(0.4),
         size=11, color=MUTED, align="left")
add_text(s, "2 / 8", Inches(0.6), Inches(6.9), Inches(12.13), Inches(0.4),
         size=11, color=MUTED, align="right")


# === Slide 3: The wedge (the killer slide) ===
s = slide_blank(prs)
set_bg(s, BG)

# Eyebrow
add_accent_bar(s, Inches(0.6), Inches(0.7), Inches(0.5))
add_text(s, "THE WEDGE", Inches(0.85), Inches(0.7), Inches(4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
         line_spacing=1.0)

# Title
add_text(s, "What Stripe physically cannot do.",
         Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.9),
         size=36, bold=True, color=TEXT, font=HEADER_FONT)

# Three stat rows
ROW_LEFT = Inches(0.6)
ROW_W = Inches(12.13)
ROW_H = Inches(1.15)
ROW_GAP = Inches(0.18)
ROW_TOP_START = Inches(2.5)

rows = [
    {
        "stat": "$0.30",
        "stat_color": MUTED,
        "label": "Stripe's minimum fee per transaction",
    },
    {
        "stat": "negative-sum",
        "stat_color": DANGER,
        "label": "Splitting one payment across 5 humans on Stripe",
    },
    {
        "stat": "~ $0.00",
        "stat_color": ACCENT,
        "label": "Splitting one payment across 5 humans on Lightning",
    },
]

for i, row in enumerate(rows):
    top = ROW_TOP_START + (ROW_H + ROW_GAP) * i
    add_rect(s, ROW_LEFT, top, ROW_W, ROW_H, CARD, corner=True, corner_radius=0.12)
    # stat (left)
    add_text(s, row["stat"], ROW_LEFT + Inches(0.4), top, Inches(5.0), ROW_H,
             size=44, bold=True, color=row["stat_color"], font=HEADER_FONT,
             anchor="middle", align="left", line_spacing=1.0)
    # label (right)
    add_text(s, row["label"], ROW_LEFT + Inches(5.5), top, Inches(6.4), ROW_H,
             size=20, color=TEXT, font=BODY_FONT, anchor="middle",
             align="left", line_spacing=1.1)

# Footer punchline
add_text(s, "We do the one thing the existing payment stack physically cannot.",
         Inches(0.6), Inches(6.85), Inches(12.13), Inches(0.5),
         size=14, italic=True, color=ACCENT_2, font=BODY_FONT, align="center")


# === Slide 4: The problem ===
s = slide_blank(prs)
set_bg(s, BG)

add_accent_bar(s, Inches(0.6), Inches(0.7), Inches(0.5))
add_text(s, "THE PROBLEM", Inches(0.85), Inches(0.7), Inches(4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
         line_spacing=1.0)

# Three-line escalating statement
add_text(s, "AI agents scale.",
         Inches(0.6), Inches(1.7), Inches(12.13), Inches(1.0),
         size=48, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.1)

add_text(s, "Model providers and platforms get paid.",
         Inches(0.6), Inches(2.85), Inches(12.13), Inches(1.0),
         size=48, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.1)

add_runs(s, [
    {"text": "The human experts whose knowledge powers the answers — ", "size": 40, "bold": True, "color": TEXT, "font": HEADER_FONT},
    {"text": "get nothing.", "size": 40, "bold": True, "color": ACCENT, "font": HEADER_FONT},
], Inches(0.6), Inches(4.05), Inches(12.13), Inches(1.8), line_spacing=1.15)

# Sub
add_text(s, "Not because nobody wants to pay them. Because the rails can't.",
         Inches(0.6), Inches(6.4), Inches(12.13), Inches(0.7),
         size=22, italic=True, color=MUTED, font=BODY_FONT, line_spacing=1.1)


# === Slide 5: How it works ===
s = slide_blank(prs)
set_bg(s, BG)

add_accent_bar(s, Inches(0.6), Inches(0.7), Inches(0.5))
add_text(s, "HOW IT WORKS", Inches(0.85), Inches(0.7), Inches(4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
         line_spacing=1.0)

# Title
add_text(s, "Live on Bitcoin mainnet.",
         Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.9),
         size=36, bold=True, color=TEXT, font=HEADER_FONT)

# 4-node horizontal flow
# 4 boxes + 3 arrows. Total width 12.13"
# Box: 2.6" wide, arrow: 0.5" wide
# 2.6*4 + 0.5*3 = 10.4 + 1.5 = 11.9. Fits in 12.13.
BOX_W = Inches(2.6)
BOX_H = Inches(2.0)
ARROW_W = Inches(0.5)
FLOW_LEFT = Inches(0.7)
FLOW_TOP = Inches(2.7)

nodes = [
    {"big": "01", "title": "Agent calls API", "sub": "GET /api/ayni/\ntributario-pe?q=..."},
    {"big": "02", "title": "402 + invoice", "sub": "Lightning invoice\n+ L402 macaroon"},
    {"big": "03", "title": "Wallet pays", "sub": "preimage verified\nsha256(preimage) == hash"},
    {"big": "04", "title": "Fan-out", "sub": "5 onward Lightning\nsends to contributors"},
]

x = FLOW_LEFT
for i, node in enumerate(nodes):
    add_rect(s, x, FLOW_TOP, BOX_W, BOX_H, CARD, corner=True, corner_radius=0.08)
    # Step number
    add_text(s, node["big"], x + Inches(0.25), FLOW_TOP + Inches(0.15), BOX_W, Inches(0.5),
             size=14, bold=True, color=ACCENT, font=HEADER_FONT, line_spacing=1.0)
    # Title
    add_text(s, node["title"], x + Inches(0.25), FLOW_TOP + Inches(0.6), BOX_W - Inches(0.5), Inches(0.6),
             size=18, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.0)
    # Sub
    add_text(s, node["sub"], x + Inches(0.25), FLOW_TOP + Inches(1.15), BOX_W - Inches(0.5), Inches(0.85),
             size=11, color=MUTED, font=MONO_FONT, line_spacing=1.2)

    # Arrow (between nodes)
    if i < len(nodes) - 1:
        ax = x + BOX_W
        ay = FLOW_TOP + BOX_H / 2 - Inches(0.18)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.05), ay, ARROW_W - Inches(0.1), Inches(0.36))
        arr.fill.solid()
        arr.fill.fore_color.rgb = ACCENT
        arr.line.fill.background()
    x = x + BOX_W + ARROW_W

# Caption
add_text(s, "End-to-end on Bitcoin mainnet.   MIT-licensed.   Public agent-skill.json manifest.",
         Inches(0.6), Inches(5.3), Inches(12.13), Inches(0.5),
         size=16, color=ACCENT_2, font=BODY_FONT, italic=True, align="center")

# Tech chips
chips = ["MoneyDevKit v0.16", "L402 protocol", "BOLT12 / Lightning Address", "Next.js 15 · Vercel"]
chip_y = Inches(6.1)
chip_h = Inches(0.5)
total_chip_w = sum(len(c) for c in chips) * 0.085 + 0.4 * len(chips)
chip_x = (13.333 - total_chip_w) / 2
chip_x_emu = Inches(chip_x)
for c in chips:
    cw = Inches(len(c) * 0.085 + 0.4)
    add_rect(s, chip_x_emu, chip_y, cw, chip_h, CARD_HI, corner=True, corner_radius=0.5)
    add_text(s, c, chip_x_emu, chip_y, cw, chip_h,
             size=11, color=TEXT, font=BODY_FONT, align="center", anchor="middle",
             line_spacing=1.0)
    chip_x_emu = chip_x_emu + cw + Inches(0.15)

add_text(s, "5 / 8", Inches(0.6), Inches(7.0), Inches(12.13), Inches(0.4),
         size=11, color=MUTED, align="right")


# === Slide 6: Live tribus ===
s = slide_blank(prs)
set_bg(s, BG)

add_accent_bar(s, Inches(0.6), Inches(0.7), Inches(0.5))
add_text(s, "LIVE TODAY", Inches(0.85), Inches(0.7), Inches(4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
         line_spacing=1.0)

add_text(s, "Two tribus shipped during the hackathon.",
         Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.9),
         size=32, bold=True, color=TEXT, font=HEADER_FONT)

# Two cards side-by-side
CARD_TOP = Inches(2.5)
CARD_HEIGHT = Inches(3.9)
CARD_GAP = Inches(0.4)
CARD_W = Inches((12.13 - 0.4) / 2)  # 5.865"

# Left card — Tributario PE
left_x = Inches(0.6)
add_rect(s, left_x, CARD_TOP, CARD_W, CARD_HEIGHT, CARD, corner=True, corner_radius=0.05)
add_text(s, "Tributario PE", left_x + Inches(0.4), CARD_TOP + Inches(0.3), CARD_W - Inches(0.8), Inches(0.6),
         size=26, bold=True, color=ACCENT, font=HEADER_FONT, line_spacing=1.0)
add_text(s, "Peruvian tax law", left_x + Inches(0.4), CARD_TOP + Inches(0.95), CARD_W - Inches(0.8), Inches(0.4),
         size=14, color=MUTED, font=BODY_FONT, italic=True)
# Big stat
add_text(s, "100", left_x + Inches(0.4), CARD_TOP + Inches(1.5), Inches(2.0), Inches(1.0),
         size=64, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.0)
add_text(s, "sat / call", left_x + Inches(2.45), CARD_TOP + Inches(1.95), Inches(2.5), Inches(0.6),
         size=18, color=MUTED, font=BODY_FONT, anchor="middle", line_spacing=1.0)

# Splits
add_text(s, "5 contributors", left_x + Inches(0.4), CARD_TOP + Inches(2.65), CARD_W - Inches(0.8), Inches(0.4),
         size=14, bold=True, color=TEXT, font=BODY_FONT)
add_text(s, "Curador 40   ·   Validador 30   ·   Contribuidor 10 × 3",
         left_x + Inches(0.4), CARD_TOP + Inches(3.05), CARD_W - Inches(0.8), Inches(0.4),
         size=14, color=ACCENT_2, font=MONO_FONT)
add_text(s, "IGV  ·  rentas 4ta / 5ta  ·  RER  ·  NRUS",
         left_x + Inches(0.4), CARD_TOP + Inches(3.45), CARD_W - Inches(0.8), Inches(0.4),
         size=12, color=MUTED, font=BODY_FONT, italic=True)

# Right card — Data Science ES
right_x = left_x + CARD_W + CARD_GAP
add_rect(s, right_x, CARD_TOP, CARD_W, CARD_HEIGHT, CARD, corner=True, corner_radius=0.05)
add_text(s, "Data Science ES", right_x + Inches(0.4), CARD_TOP + Inches(0.3), CARD_W - Inches(0.8), Inches(0.6),
         size=26, bold=True, color=ACCENT, font=HEADER_FONT, line_spacing=1.0)
add_text(s, "Spanish / LatAm DS practice", right_x + Inches(0.4), CARD_TOP + Inches(0.95), CARD_W - Inches(0.8), Inches(0.4),
         size=14, color=MUTED, font=BODY_FONT, italic=True)
add_text(s, "75", right_x + Inches(0.4), CARD_TOP + Inches(1.5), Inches(2.0), Inches(1.0),
         size=64, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.0)
add_text(s, "sat / call", right_x + Inches(2.05), CARD_TOP + Inches(1.95), Inches(2.5), Inches(0.6),
         size=18, color=MUTED, font=BODY_FONT, anchor="middle", line_spacing=1.0)

add_text(s, "3 contributors", right_x + Inches(0.4), CARD_TOP + Inches(2.65), CARD_W - Inches(0.8), Inches(0.4),
         size=14, bold=True, color=TEXT, font=BODY_FONT)
add_text(s, "Curador 50   ·   Contribuidor 25 × 2",
         right_x + Inches(0.4), CARD_TOP + Inches(3.05), CARD_W - Inches(0.8), Inches(0.4),
         size=14, color=ACCENT_2, font=MONO_FONT)
add_text(s, "models  ·  datasets  ·  regional best practices",
         right_x + Inches(0.4), CARD_TOP + Inches(3.45), CARD_W - Inches(0.8), Inches(0.4),
         size=12, color=MUTED, font=BODY_FONT, italic=True)

# Footer
add_text(s, "Adding a tribu = config change in data/tribus.json. Not a code change.",
         Inches(0.6), Inches(6.7), Inches(12.13), Inches(0.5),
         size=14, italic=True, color=ACCENT_2, font=BODY_FONT, align="center")


# === Slide 7: Why now / Why us ===
s = slide_blank(prs)
set_bg(s, BG)

add_accent_bar(s, Inches(0.6), Inches(0.7), Inches(0.5))
add_text(s, "TIMING & TEAM", Inches(0.85), Inches(0.7), Inches(4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
         line_spacing=1.0)

# Two columns
COL_TOP = Inches(1.9)
COL_W = Inches((12.13 - 0.4) / 2)
COL_GAP = Inches(0.4)

# Why now
left_x = Inches(0.6)
add_text(s, "Why now", left_x, COL_TOP, COL_W, Inches(0.7),
         size=32, bold=True, color=TEXT, font=HEADER_FONT)

why_now = [
    "L402 finally ergonomic   ·   MoneyDevKit v0.16",
    "AI agents went production in 2026",
    "Lightning capacity mature",
    "Stablecoin rails replace one gatekeeper with another. Lightning doesn't.",
]
y = COL_TOP + Inches(0.95)
for i, line in enumerate(why_now):
    # bullet dot
    add_rect(s, left_x, y + Inches(0.18), Inches(0.12), Inches(0.12), ACCENT, corner=True, corner_radius=0.5)
    add_text(s, line, left_x + Inches(0.3), y, COL_W - Inches(0.3), Inches(0.7),
             size=18, color=TEXT, font=BODY_FONT, line_spacing=1.2)
    y = y + Inches(0.85)

# Why us
right_x = left_x + COL_W + COL_GAP
add_text(s, "Why us", right_x, COL_TOP, COL_W, Inches(0.7),
         size=32, bold=True, color=TEXT, font=HEADER_FONT)

why_us = [
    "4 builders in Lima, Peru",
    "The communities we monetize are people we can call by name",
    "Shipped end-to-end on Bitcoin mainnet — in a weekend",
    "MoneyDevKit accessible in Spanish via Discord — natural fit",
]
y = COL_TOP + Inches(0.95)
for line in why_us:
    add_rect(s, right_x, y + Inches(0.18), Inches(0.12), Inches(0.12), ACCENT, corner=True, corner_radius=0.5)
    add_text(s, line, right_x + Inches(0.3), y, COL_W - Inches(0.3), Inches(0.7),
             size=18, color=TEXT, font=BODY_FONT, line_spacing=1.2)
    y = y + Inches(0.85)

add_text(s, "7 / 8", Inches(0.6), Inches(7.0), Inches(12.13), Inches(0.4),
         size=11, color=MUTED, align="right")


# === Slide 8: The ask ===
s = slide_blank(prs)
set_bg(s, BG)

# Big accent block on right edge
add_rect(s, Inches(12.5), Inches(0), Inches(0.85), Inches(7.5), ACCENT)

add_accent_bar(s, Inches(0.6), Inches(0.7), Inches(0.5))
add_text(s, "THE ASK", Inches(0.85), Inches(0.7), Inches(4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
         line_spacing=1.0)

# Big centered statement
add_text(s, "Spiral Challenge 02.",
         Inches(0.6), Inches(2.0), Inches(11.3), Inches(1.3),
         size=64, bold=True, color=TEXT, font=HEADER_FONT, align="center", line_spacing=1.05)

add_text(s, "Back the only payment rail that can.",
         Inches(0.6), Inches(3.5), Inches(11.3), Inches(1.3),
         size=56, bold=True, color=ACCENT, font=HEADER_FONT, align="center", line_spacing=1.05)

# Stop-talking line
add_text(s, "Ayni is reciprocity, encoded in software, at sub-cent scale.",
         Inches(0.6), Inches(5.1), Inches(11.3), Inches(0.6),
         size=20, italic=True, color=ACCENT_2, font=BODY_FONT, align="center", line_spacing=1.1)

# URLs
add_text(s, "ayniw.com     ·     github.com/d3nn1sVZ/Ayni-agents",
         Inches(0.6), Inches(6.6), Inches(11.3), Inches(0.5),
         size=18, color=TEXT, font=MONO_FONT, align="center")


# --- Save ---
out = Path("ayni-pitch/Ayni-pitch-deck.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"Wrote {out}  ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")
