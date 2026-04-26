"""Build the Ayni 60-second TECH video deck via python-pptx.

6 slides, screen-record alongside narration in pitch-tech-1min.md.
Output: ayni-pitch/Ayni-tech-deck.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette (same as main deck for consistency)
BG       = RGBColor(0x0B, 0x10, 0x20)
CARD     = RGBColor(0x18, 0x21, 0x38)
CARD_HI  = RGBColor(0x22, 0x2D, 0x4A)
ACCENT   = RGBColor(0xF7, 0x93, 0x1A)
ACCENT_2 = RGBColor(0xFF, 0xD5, 0x80)
TEXT     = RGBColor(0xEC, 0xEC, 0xF0)
MUTED    = RGBColor(0x8B, 0x93, 0xA7)
GREEN    = RGBColor(0x6E, 0xCC, 0x9A)

HEADER_FONT = "Calibri"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, *, line=None, corner=False, corner_radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    if corner and corner_radius is not None:
        shp.adjustments[0] = corner_radius
    return shp


def add_text(slide, text, left, top, width, height, *, size=18, bold=False, italic=False,
             color=None, align="left", font=BODY_FONT, anchor="top", line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tb


def add_accent_bar(slide, left, top, height, *, width=Inches(0.07), color=ACCENT):
    return add_rect(slide, left, top, width, height, color)


def add_arrow(slide, left, top, width, height, color=ACCENT):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_timing_badge(slide, text):
    add_rect(slide, Inches(11.4), Inches(0.4), Inches(1.5), Inches(0.45),
             CARD_HI, corner=True, corner_radius=0.5)
    add_text(slide, text, Inches(11.4), Inches(0.4), Inches(1.5), Inches(0.45),
             size=12, color=TEXT, font=MONO_FONT, align="center", anchor="middle",
             line_spacing=1.0)


def add_eyebrow(slide, text):
    add_accent_bar(slide, Inches(0.6), Inches(0.7), Inches(0.5))
    add_text(slide, text, Inches(0.85), Inches(0.7), Inches(6), Inches(0.5),
             size=14, bold=True, color=ACCENT, font=HEADER_FONT, anchor="middle",
             line_spacing=1.0)


def add_notes(slide, text):
    """Add speaker notes to a slide. Visible in Presenter View only."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# ---------------- BUILD ----------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# === Slide 1: Hero card (0:00 – 0:08) ===
s = slide_blank(prs)
set_bg(s, BG)
add_timing_badge(s, "0:00 – 0:08")

add_text(s, "ayni", Inches(0), Inches(1.4), SLIDE_W, Inches(2.6),
         size=160, bold=True, color=ACCENT, font=HEADER_FONT, align="center",
         anchor="middle", line_spacing=1.0)

add_text(s, "agent-paid expert APIs   ·   over Lightning",
         Inches(0), Inches(4.1), SLIDE_W, Inches(0.8),
         size=32, bold=True, color=TEXT, font=HEADER_FONT, align="center",
         line_spacing=1.0)

add_text(s, "AI agents pay sub-cent in Lightning. Payment auto-splits to the contributors who maintain the knowledge.",
         Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.2),
         size=18, color=MUTED, font=BODY_FONT, align="center", italic=True,
         line_spacing=1.3)

add_text(s, "Hack Nation 5 — Tech Video — 60 sec",
         Inches(0), Inches(6.8), SLIDE_W, Inches(0.4),
         size=12, color=MUTED, font=BODY_FONT, align="center")

add_notes(s, """[0:00 – 0:08]   HERO

Ayni: AI agents pay sub-cent in Lightning, and every payment auto-splits to the contributors who maintain the knowledge.

— beat — let it land before clicking to slide 2.""")


# === Slide 2: System Architecture (0:08 – 0:25) ===
s = slide_blank(prs)
set_bg(s, BG)
add_timing_badge(s, "0:08 – 0:25")
add_eyebrow(s, "SYSTEM ARCHITECTURE")

add_text(s, "How it fits together.",
         Inches(0.6), Inches(1.35), Inches(12.13), Inches(0.9),
         size=34, bold=True, color=TEXT, font=HEADER_FONT)

# --- Diagram canvas ---
DIAG_TOP = Inches(2.55)
DIAG_H = Inches(4.3)

# === Column 1: Agent ===
AGENT_X = Inches(0.6)
AGENT_W = Inches(2.3)
agent_box_y = DIAG_TOP + Inches(0.95)
agent_box_h = Inches(2.4)
add_rect(s, AGENT_X, agent_box_y, AGENT_W, agent_box_h, CARD, corner=True, corner_radius=0.06)
add_text(s, "AI Agent", AGENT_X, agent_box_y + Inches(0.35), AGENT_W, Inches(0.6),
         size=22, bold=True, color=ACCENT, font=HEADER_FONT, align="center", line_spacing=1.0)
add_text(s, "+ wallet", AGENT_X, agent_box_y + Inches(1.0), AGENT_W, Inches(0.5),
         size=16, color=TEXT, font=HEADER_FONT, align="center", line_spacing=1.0)
add_text(s, "MDK · Phoenix\nAlby · Strike", AGENT_X, agent_box_y + Inches(1.55), AGENT_W, Inches(0.7),
         size=10, color=MUTED, font=MONO_FONT, align="center", line_spacing=1.3)

# === Column 2: Ayni system ===
AYNI_X = Inches(3.5)
AYNI_W = Inches(5.9)
add_rect(s, AYNI_X, DIAG_TOP, AYNI_W, DIAG_H, CARD, corner=True, corner_radius=0.04, line=ACCENT)

# Title bar inside
add_text(s, "Ayni  —  Next.js 15 on Vercel",
         AYNI_X + Inches(0.2), DIAG_TOP + Inches(0.15), AYNI_W - Inches(0.4), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, align="center", line_spacing=1.0,
         anchor="middle")

# 2x2 sub-grid
SUB_TOP = DIAG_TOP + Inches(0.75)
SUB_INNER_X = AYNI_X + Inches(0.25)
SUB_INNER_W = AYNI_W - Inches(0.5)
SUB_GAP = Inches(0.2)
CELL_W = Inches((5.9 - 0.5 - 0.2) / 2)  # ~2.6"
SUB_INNER_H = DIAG_H - Inches(0.95)
CELL_H = Inches((4.3 - 0.95 - 0.2) / 2)  # ~1.575"

sub_boxes = [
    {"row": 0, "col": 0, "title": "manifest", "tag": "discovery",
     "code": "/.well-known/\nagent-skill.json"},
    {"row": 0, "col": 1, "title": "API route", "tag": "L402 paywall",
     "code": "/api/ayni/[plugin]\nMDK · withPayment"},
    {"row": 1, "col": 0, "title": "splitter", "tag": "fan-out",
     "code": "scripts/demo-flow.sh\nfloor-then-residual"},
    {"row": 1, "col": 1, "title": "dashboard", "tag": "live UI",
     "code": "/api/payouts/stream\nSSE  ·  2-phase"},
]

for box in sub_boxes:
    bx = SUB_INNER_X + (CELL_W + SUB_GAP) * box["col"]
    by = SUB_TOP + (CELL_H + SUB_GAP) * box["row"]
    add_rect(s, bx, by, CELL_W, CELL_H, CARD_HI, corner=True, corner_radius=0.06)
    add_text(s, box["title"], bx + Inches(0.2), by + Inches(0.15),
             CELL_W - Inches(0.4), Inches(0.4),
             size=15, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.0)
    add_text(s, box["tag"], bx + Inches(0.2), by + Inches(0.55),
             CELL_W - Inches(0.4), Inches(0.3),
             size=10, italic=True, color=ACCENT, font=BODY_FONT, line_spacing=1.0)
    add_text(s, box["code"], bx + Inches(0.2), by + Inches(0.85),
             CELL_W - Inches(0.4), CELL_H - Inches(0.95),
             size=10, color=ACCENT_2, font=MONO_FONT, line_spacing=1.3)

# === Column 3: Lightning + Contributors stacked ===
RIGHT_X = Inches(10.0)
RIGHT_W = Inches(2.7)
LN_TOP = DIAG_TOP
LN_H = Inches(1.95)
add_rect(s, RIGHT_X, LN_TOP, RIGHT_W, LN_H, CARD, corner=True, corner_radius=0.06)
add_text(s, "Lightning", RIGHT_X, LN_TOP + Inches(0.25), RIGHT_W, Inches(0.5),
         size=20, bold=True, color=ACCENT, font=HEADER_FONT, align="center", line_spacing=1.0)
add_text(s, "Bitcoin mainnet", RIGHT_X, LN_TOP + Inches(0.8), RIGHT_W, Inches(0.4),
         size=14, color=TEXT, font=HEADER_FONT, align="center", line_spacing=1.0)
add_text(s, "BOLT11 invoices", RIGHT_X, LN_TOP + Inches(1.25), RIGHT_W, Inches(0.4),
         size=11, color=MUTED, font=MONO_FONT, align="center", line_spacing=1.0)

CW_TOP = LN_TOP + LN_H + Inches(0.25)
CW_H = DIAG_TOP + DIAG_H - CW_TOP
add_rect(s, RIGHT_X, CW_TOP, RIGHT_W, CW_H, CARD, corner=True, corner_radius=0.06)
add_text(s, "Contributors", RIGHT_X, CW_TOP + Inches(0.2), RIGHT_W, Inches(0.5),
         size=18, bold=True, color=ACCENT, font=HEADER_FONT, align="center", line_spacing=1.0)
add_text(s, "BOLT12  /  LN-Address", RIGHT_X, CW_TOP + Inches(0.75), RIGHT_W, Inches(0.4),
         size=11, color=TEXT, font=MONO_FONT, align="center", line_spacing=1.0)

# Five contributor wallet circles
mini_y = CW_TOP + Inches(1.35)
mini_d = Inches(0.34)
mini_gap = Inches(0.13)
total_mini_w = mini_d * 5 + mini_gap * 4
mini_start_x = RIGHT_X + (RIGHT_W - total_mini_w) / 2
for i in range(5):
    cx = mini_start_x + (mini_d + mini_gap) * i
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, mini_y, mini_d, mini_d)
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT_2
    circ.line.fill.background()

# === Arrows ===
# Arrow 1: Agent -> Ayni (right-pointing, mid-height)
ARROW_W = Inches(0.5)
arr1_y = DIAG_TOP + Inches(2.0)
add_arrow(s, AGENT_X + AGENT_W + Inches(0.05), arr1_y, ARROW_W - Inches(0.1), Inches(0.32))
add_text(s, "1.  call",
         AGENT_X + AGENT_W, arr1_y - Inches(0.4),
         ARROW_W, Inches(0.3),
         size=10, bold=True, color=TEXT, font=HEADER_FONT, align="center", line_spacing=1.0)
add_text(s, "← answer",
         AGENT_X + AGENT_W, arr1_y + Inches(0.35),
         ARROW_W, Inches(0.3),
         size=9, color=MUTED, font=BODY_FONT, align="center", line_spacing=1.0)

# Arrow 2: Ayni -> Lightning (top right zone)
arr2_y = LN_TOP + Inches(0.85)
add_arrow(s, AYNI_X + AYNI_W + Inches(0.05), arr2_y, ARROW_W - Inches(0.1), Inches(0.32))
add_text(s, "2.  invoice",
         AYNI_X + AYNI_W, arr2_y - Inches(0.4),
         ARROW_W, Inches(0.3),
         size=9, bold=True, color=TEXT, font=HEADER_FONT, align="center", line_spacing=1.0)
add_text(s, "← preimage",
         AYNI_X + AYNI_W, arr2_y + Inches(0.35),
         ARROW_W, Inches(0.3),
         size=8, color=MUTED, font=BODY_FONT, align="center", line_spacing=1.0)

# Arrow 3: Ayni splitter -> Contributors (bottom right zone)
arr3_y = CW_TOP + Inches(0.95)
add_arrow(s, AYNI_X + AYNI_W + Inches(0.05), arr3_y, ARROW_W - Inches(0.1), Inches(0.32))
add_text(s, "3.  fan-out",
         AYNI_X + AYNI_W, arr3_y - Inches(0.4),
         ARROW_W, Inches(0.3),
         size=9, bold=True, color=TEXT, font=HEADER_FONT, align="center", line_spacing=1.0)

# Footer caption
add_text(s, "(1) agent calls API   →   (2) MDK issues + verifies Lightning invoice   →   (3) splitter fans out to all contributors",
         Inches(0.6), Inches(7.0), Inches(12.13), Inches(0.4),
         size=11, italic=True, color=ACCENT_2, font=BODY_FONT, align="center", line_spacing=1.0)

add_notes(s, """[0:08 – 0:25]   SYSTEM ARCHITECTURE   ★ centerpiece slide ★

(point at the AGENT column on the left)
Stack: Next.js 15 on Vercel, MoneyDevKit 0.16, L402 over Bitcoin mainnet.

(sweep right, into the AYNI middle container)
An agent discovers our endpoints through a public agent-skill manifest, then calls the API.

(point at arrow 2 — Ayni → Lightning)
The API returns a 402 with a Lightning invoice and a macaroon. The agent's wallet pays. We verify the preimage and return the answer.

(point at arrow 3 — Splitter → Contributors)
Then the splitter fans out real Lightning sends to each contributor's BOLT12 or Lightning Address.

— 17 seconds. Densest slide. Practice this one alone first. Use cursor / laser pointer (Ctrl+L in slideshow mode) to point at each zone in order.""")


# === Slide 3: L402 zoom (0:25 – 0:35) ===
s = slide_blank(prs)
set_bg(s, BG)
add_timing_badge(s, "0:25 – 0:35")
add_eyebrow(s, "L402 ZOOM")

add_text(s, "Inside the API route.",
         Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.9),
         size=36, bold=True, color=TEXT, font=HEADER_FONT)

# 4-step flow with arrows
nodes = [
    {"big": "01", "title": "Agent calls", "code": "GET /api/ayni/\ntributario-pe?q=igv"},
    {"big": "02", "title": "402 returned", "code": "Lightning invoice\n+ L402 macaroon"},
    {"big": "03", "title": "Pay + retry", "code": "Authorization:\nL402 <mac>:<preimage>"},
    {"big": "04", "title": "Verify + 200", "code": "sha256(preimage)\n=== paymentHash"},
]

BOX_W = Inches(2.6)
BOX_H = Inches(2.4)
ARROW_W = Inches(0.5)
FLOW_LEFT = Inches(0.7)
FLOW_TOP = Inches(2.6)

x = FLOW_LEFT
for i, node in enumerate(nodes):
    add_rect(s, x, FLOW_TOP, BOX_W, BOX_H, CARD, corner=True, corner_radius=0.08)
    add_text(s, node["big"], x + Inches(0.3), FLOW_TOP + Inches(0.2), BOX_W, Inches(0.5),
             size=14, bold=True, color=ACCENT, font=HEADER_FONT, line_spacing=1.0)
    add_text(s, node["title"], x + Inches(0.3), FLOW_TOP + Inches(0.7), BOX_W - Inches(0.6), Inches(0.6),
             size=18, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.0)
    add_text(s, node["code"], x + Inches(0.3), FLOW_TOP + Inches(1.3), BOX_W - Inches(0.6), Inches(1.0),
             size=11, color=ACCENT_2, font=MONO_FONT, line_spacing=1.25)
    if i < len(nodes) - 1:
        ax = x + BOX_W
        ay = FLOW_TOP + BOX_H / 2 - Inches(0.18)
        add_arrow(s, ax + Inches(0.05), ay, ARROW_W - Inches(0.1), Inches(0.36))
    x = x + BOX_W + ARROW_W

add_text(s, "End-to-end on Bitcoin mainnet. Real invoices. Real preimage verification.",
         Inches(0.6), Inches(5.5), Inches(12.13), Inches(0.5),
         size=16, italic=True, color=ACCENT_2, font=BODY_FONT, align="center")

add_notes(s, """[0:25 – 0:35]   L402 ZOOM

Inside the API route: MDK returns the 402 with the Lightning invoice and macaroon. The agent retries with the preimage. We verify SHA-256 against the payment hash before answering.

— stress: "402", "macaroon", "preimage", "SHA-256". 10 seconds.""")


# === Slide 4: Splitter (0:35 – 0:46) ===
s = slide_blank(prs)
set_bg(s, BG)
add_timing_badge(s, "0:35 – 0:46")
add_eyebrow(s, "ONWARD FAN-OUT")

add_text(s, "Splitter: floor-then-residual.",
         Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.9),
         size=36, bold=True, color=TEXT, font=HEADER_FONT)

PSEUDO_LEFT = Inches(0.6)
PSEUDO_TOP = Inches(2.6)
PSEUDO_W = Inches(7.0)
PSEUDO_H = Inches(3.7)
add_rect(s, PSEUDO_LEFT, PSEUDO_TOP, PSEUDO_W, PSEUDO_H, CARD, corner=True, corner_radius=0.04)

pseudo = (
    "# floor-then-residual\n"
    "amounts = [floor(total * pct) for pct in splits]\n"
    "residual = total - sum(amounts)\n"
    "amounts[0] += residual    # curador absorbs drift\n"
    "\n"
    "assert sum(amounts) == total\n"
    "for wallet, sats in zip(wallets, amounts):\n"
    "    pay_lightning(wallet, sats)   # BOLT12 / LN-Address"
)
add_text(s, pseudo, PSEUDO_LEFT + Inches(0.4), PSEUDO_TOP + Inches(0.35),
         PSEUDO_W - Inches(0.8), PSEUDO_H - Inches(0.7),
         size=16, color=ACCENT_2, font=MONO_FONT, line_spacing=1.45)

EX_LEFT = PSEUDO_LEFT + PSEUDO_W + Inches(0.4)
EX_W = Inches(12.73 - 7.4)
add_rect(s, EX_LEFT, PSEUDO_TOP, EX_W, PSEUDO_H, CARD_HI, corner=True, corner_radius=0.04)

add_text(s, "Example", EX_LEFT + Inches(0.3), PSEUDO_TOP + Inches(0.25), EX_W - Inches(0.6), Inches(0.5),
         size=14, bold=True, color=ACCENT, font=HEADER_FONT, line_spacing=1.0)
add_text(s, "100 sat → Tributario PE", EX_LEFT + Inches(0.3), PSEUDO_TOP + Inches(0.7),
         EX_W - Inches(0.6), Inches(0.5),
         size=18, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.0)

example_lines = [("40", "Curador"), ("30", "Validador"),
                 ("10", "Contribuidor"), ("10", "Contribuidor"), ("10", "Contribuidor")]
ey = PSEUDO_TOP + Inches(1.45)
for sats, role in example_lines:
    add_text(s, sats, EX_LEFT + Inches(0.3), ey, Inches(1.2), Inches(0.4),
             size=20, bold=True, color=ACCENT, font=MONO_FONT, line_spacing=1.0)
    add_text(s, role, EX_LEFT + Inches(1.5), ey, EX_W - Inches(2.0), Inches(0.4),
             size=16, color=TEXT, font=BODY_FONT, line_spacing=1.0, anchor="middle")
    ey = ey + Inches(0.42)

add_text(s, "Σ = 100 sat. No drift.",
         EX_LEFT + Inches(0.3), PSEUDO_TOP + Inches(3.1), EX_W - Inches(0.6), Inches(0.4),
         size=14, italic=True, color=GREEN, font=MONO_FONT, line_spacing=1.0)

add_text(s, "Each wallet paid by BOLT12 or Lightning Address via the MDK agent-wallet.",
         Inches(0.6), Inches(6.65), Inches(12.13), Inches(0.5),
         size=14, italic=True, color=MUTED, font=BODY_FONT, align="center")

add_notes(s, """[0:35 – 0:46]   ONWARD FAN-OUT

The splitter uses floor-then-residual rounding — per-wallet sats sum exactly to the total. A 100-sat call to Tributario PE pays 40-30-10-10-10 to five wallets in seconds.

— stress: "floor-then-residual", "exactly", "40-30-10-10-10". 11 seconds.""")


# === Slide 5: Critical fix (0:46 – 0:55) ===
s = slide_blank(prs)
set_bg(s, BG)
add_timing_badge(s, "0:46 – 0:55")
add_eyebrow(s, "IMPLEMENTATION")

add_text(s, "The fix that unlocked end-to-end.",
         Inches(0.6), Inches(1.3), Inches(12.13), Inches(0.9),
         size=34, bold=True, color=TEXT, font=HEADER_FONT)

THREE_TOP = Inches(2.55)
THREE_H = Inches(3.5)
THREE_W = Inches((12.13 - 0.6) / 3)
THREE_GAP = Inches(0.3)

cards = [
    {"label": "SYMPTOM", "title": "L402 route 500'd at runtime",
     "body": "ws.mask() broke when bundled by webpack. Native lightning-js .node binding can't be webpacked at all."},
    {"label": "CAUSE", "title": "Webpack vs. native bindings",
     "body": "Next.js bundled both modules into the route. Native code refused to load; ws was mangled by minification."},
    {"label": "FIX", "title": "Externalize + Node runtime",
     "body": "next.config.mjs externals: ['ws','@moneydevkit/lightning-js']. Route handler: export const runtime = 'nodejs'."},
]

for i, c in enumerate(cards):
    x = Inches(0.6) + (THREE_W + THREE_GAP) * i
    add_rect(s, x, THREE_TOP, THREE_W, THREE_H, CARD, corner=True, corner_radius=0.05)
    add_rect(s, x + Inches(0.3), THREE_TOP + Inches(0.3), Inches(1.4), Inches(0.4),
             CARD_HI, corner=True, corner_radius=0.5)
    add_text(s, c["label"], x + Inches(0.3), THREE_TOP + Inches(0.3), Inches(1.4), Inches(0.4),
             size=11, bold=True, color=ACCENT, font=HEADER_FONT, align="center", anchor="middle",
             line_spacing=1.0)
    add_text(s, c["title"], x + Inches(0.3), THREE_TOP + Inches(0.85), THREE_W - Inches(0.6), Inches(1.0),
             size=18, bold=True, color=TEXT, font=HEADER_FONT, line_spacing=1.15)
    add_text(s, c["body"], x + Inches(0.3), THREE_TOP + Inches(2.0), THREE_W - Inches(0.6), Inches(1.4),
             size=13, color=MUTED, font=BODY_FONT, line_spacing=1.3)

add_text(s, "commit  09efcb6",
         Inches(0.6), Inches(6.55), Inches(12.13), Inches(0.4),
         size=14, color=ACCENT, font=MONO_FONT, align="center", bold=True, line_spacing=1.0)
add_text(s, "fix(l402): externalize ws + lightning native module, set nodejs runtime",
         Inches(0.6), Inches(6.95), Inches(12.13), Inches(0.4),
         size=12, italic=True, color=MUTED, font=MONO_FONT, align="center", line_spacing=1.0)

add_notes(s, """[0:46 – 0:55]   IMPLEMENTATION — THE FIX

Critical fix: the WebSocket lib and the native Lightning binding broke under webpack. We externalized them and pinned the route to the Node.js runtime. That commit unlocked end-to-end.

— stress: "externalized", "Node.js runtime", "unlocked end-to-end". 9 seconds.""")


# === Slide 6: Closing card (0:55 – 1:00) ===
s = slide_blank(prs)
set_bg(s, BG)
add_timing_badge(s, "0:55 – 1:00")

add_text(s, "Live on mainnet.",
         Inches(0), Inches(2.1), SLIDE_W, Inches(1.3),
         size=72, bold=True, color=TEXT, font=HEADER_FONT, align="center", line_spacing=1.05)

add_text(s, "MIT on GitHub.",
         Inches(0), Inches(3.5), SLIDE_W, Inches(1.3),
         size=72, bold=True, color=ACCENT, font=HEADER_FONT, align="center", line_spacing=1.05)

add_text(s, "ayniw.com",
         Inches(0), Inches(5.6), SLIDE_W, Inches(0.6),
         size=28, color=TEXT, font=MONO_FONT, align="center", line_spacing=1.0)
add_text(s, "github.com/d3nn1sVZ/Ayni-agents",
         Inches(0), Inches(6.2), SLIDE_W, Inches(0.6),
         size=24, color=MUTED, font=MONO_FONT, align="center", line_spacing=1.0)

add_notes(s, """[0:55 – 1:00]   CLOSING

Live on mainnet at ayniw.com. MIT on GitHub.

— stop talking. Don't read the URL twice. 5 seconds.""")


out = Path("ayni-pitch/Ayni-tech-deck.pptx")
prs.save(str(out))
print(f"Wrote {out}  ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")
